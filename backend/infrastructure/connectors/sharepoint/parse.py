"""Binary file content -> plain text, dispatched by filename extension.

Graph's ``/content`` endpoint returns raw bytes regardless of file type — a
``.docx``/``.pdf``/``.xlsx`` decoded as UTF-8 is garbage, so ``sharepoint_read``
needs real parsing before handing content to the agent.
"""

from __future__ import annotations

import io
import logging

logger = logging.getLogger(__name__)


def _parse_docx(content: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pdf(content: bytes) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(f"[Page {i + 1}]\n{text}" for i, text in enumerate(pages) if text.strip())


def _parse_xlsx(content: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


_PARSERS = {
    "docx": _parse_docx,
    "pdf": _parse_pdf,
    "xlsx": _parse_xlsx,
    "pptx": _parse_pptx,
}

_TEXT_EXTENSIONS = {"txt", "md", "csv", "json", "yaml", "yml", "log"}


def extract_text(filename: str, content: bytes) -> str:
    """Best-effort text extraction. Unknown/unparseable content falls back to
    a note instead of raw bytes garbage — never silently returns binary."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in _TEXT_EXTENSIONS:
        return content.decode("utf-8", errors="replace")

    parser = _PARSERS.get(ext)
    if parser is None:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return f"[binary content, {len(content)} bytes — unsupported file type {ext or '(no extension)'!r}, cannot extract text]"

    try:
        text = parser(content)
    except Exception as exc:
        logger.warning("extract_text: failed to parse %r (%s): %s", filename, ext, exc)
        return f"[error extracting text from {filename!r} ({ext}): {exc}]"

    return text or f"[{filename} parsed successfully but contains no extractable text]"
